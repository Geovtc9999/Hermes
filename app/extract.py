"""Extraction de texte selon le type de fichier + dérivation version/domaine."""
from __future__ import annotations

import io
import os

# Types texte traités en v1. (.chm, .doc binaire, .zip, médias : ignorés en v1.)
SUPPORTED = {".pdf", ".docx", ".pptx", ".txt", ".md", ".csv", ".html", ".htm", ".xml"}


def supported(key: str) -> bool:
    return os.path.splitext(key)[1].lower() in SUPPORTED


def derive_meta(key: str) -> tuple[str | None, str | None]:
    """À partir de la clé bucket (organisée par préfixe), déduit (version, domaine)."""
    parts = key.split("/")
    top = parts[0] if parts else ""
    version_map = {
        "v8": "V8", "v10": "V10", "v11": "V11",
        "produit-chm-2018": "CHM2018", "version-2010": "V2010",
    }
    version = version_map.get(top)
    if top.startswith("core-y2"):
        version = "CORE-Y2"
    domaine = parts[1] if len(parts) > 1 else top
    return version, domaine


def extract_text(key: str, data: bytes) -> str:
    ext = os.path.splitext(key)[1].lower()
    try:
        if ext == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(data))
            return "\n".join((p.extract_text() or "") for p in reader.pages)
        if ext == ".docx":
            import docx
            d = docx.Document(io.BytesIO(data))
            return "\n".join(p.text for p in d.paragraphs)
        if ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data))
            chunks = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        chunks.append(shape.text_frame.text)
            return "\n".join(chunks)
        if ext in (".html", ".htm", ".xml"):
            from bs4 import BeautifulSoup
            return BeautifulSoup(data, "html.parser").get_text(" ", strip=True)
        if ext in (".txt", ".md", ".csv"):
            return data.decode("utf-8", errors="replace")
    except Exception as e:  # extraction best-effort
        return f""  # signalée vide -> sera ignorée par l'ingestion
    return ""
